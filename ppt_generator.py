from pptx import Presentation
import os
import utils as u  # 引入 utils 模組
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def adjust_bullet_points(prs, font_size=18, line_spacing=0.6):
    """調整條列摘要的字體大小和排版"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.startswith("-") or paragraph.text.startswith("•"):
                        paragraph.font.size = Pt(font_size)
                        paragraph.line_spacing = line_spacing
                        paragraph.alignment = PP_ALIGN.LEFT  # 左對齊
                        paragraph.space_after = Pt(5)  # 設置段後距離
                        paragraph.font.bold = False  # 正文字體不加粗
    return prs

def remove_empty_placeholders(prs):
    """移除所有幻燈片上的空白占位符（如 '按一下以新增標題'）"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text.strip():  # 如果文本為空則刪除該占位符
                sp = shape
                slide.shapes._spTree.remove(sp._element)
    return prs

def create_presentation(pdf_path, image_folder):
    """
    根據 PDF 和圖片生成 PPT，包含摘要、詳細整理與圖片解釋。返回生成的 PPT 文件名。
    """
    prs = Presentation()

    pdf_text = u.read_pdf(pdf_path)
    title = u.get_title(pdf_text)

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    u.adjust_title_font(slide, title or "未能提取標題")

    summary, details = u.generate_pdf_summary_and_details(pdf_text)

    for i, chunk in enumerate(u.split_text_by_points(summary, points_per_slide=2)):
        u.add_text_slide(prs, f"PDF 條列摘要（第 {i + 1} 部分）", "\n".join(chunk), bullet=True, font_size=18)

    for i, chunk in enumerate(u.split_text_by_points(details, points_per_slide=2)):
        u.add_text_slide(prs, f"PDF 詳細整理（第 {i + 1} 部分）", "\n".join(chunk), bullet=False, font_size=18)

    adjust_bullet_points(prs, font_size=18, line_spacing=0.6)

    images = sorted([
        os.path.join(image_folder, img)
        for img in os.listdir(image_folder)
        if img.lower().endswith(('.png', '.jpg', '.jpeg'))
    ])
    for i, image_path in enumerate(images):
        description = u.generate_image_description(image_path)
        u.add_image_and_description_to_slide(prs, image_path, description)

    # 移除空白占位符並調整條列格式
    prs = remove_empty_placeholders(prs)
    prs = adjust_bullet_points(prs, font_size=20, line_spacing=1.2)

    output_path = 'pdf_and_images_analysis.pptx'
    prs.save(output_path)
    return output_path