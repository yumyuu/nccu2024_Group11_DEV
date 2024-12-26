from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
# from PyPDF2 import PdfReader
import fitz
import google.generativeai as genai
import os
import base64
import uuid

genai.configure(api_key='AIzaSyADv9y5ye8btqr12Wlwo7FD-pPJBUMMc_A')
# 到 https://ai.google.dev/gemini-api/docs/api-key 申請 API Key
model = genai.GenerativeModel(model_name='gemini-1.5-flash') # 選擇模型

def read_pdf(pdf_path):
    try:
        with fitz.open(pdf_path) as doc:  # 兼容新版 fitz
            return "\n".join(page.get_text() for page in doc)
    except Exception as e:
        print(f"讀取 PDF 發生錯誤: {e}")
        return ""

def get_unique_filename(filename):
    """為文件生成唯一名稱"""
    name, ext = os.path.splitext(filename)
    unique_name = f"{name}_{uuid.uuid4().hex[:8]}{ext}"
    return unique_name

# Google Generative AI 調用函數
def call_generative_ai(prompt):
    """通用的 Google Generative AI 調用函數"""
    try:
        response = model.generate_content(prompt)
        # print(f"API 返回值: {response}") # 測試用
        return response
    except Exception as e:
        print(f"AI 調用失敗：{e}")
        return None
    
# 從 API 返回值中提取內容
def extract_text(response):
    try:
        candidates = response.candidates
        if candidates and len(candidates) > 0:
            parts = candidates[0].content.parts
            if parts and len(parts) > 0:
                return parts[0].text  # 獲取文本
        return "無法提取內容，返回結構不完整"
    except (KeyError, IndexError, AttributeError, TypeError) as e:
        print(f"提取 text 時發生錯誤：{e}")
        return "無法提取內容"


def get_title(text, top_text=500):
    """生成論文標題"""
    try:
        prompt = f"只給我這篇論文的標題：\n\n{text[:top_text]}"
        response = model.generate_content(prompt)

        # print(f"API 返回類型: {type(response)}") # 測試用
        # print(f"API 返回內容: {response}")
        # print("返回對象的類型:", type(response))
        # print("返回對象的方法和屬性:", dir(response))
        # print("候選者內容:", response.candidates[0].content)
        # print("候選者屬性:", dir(response.candidates[0].content))

        return extract_text(response) or "未能生成標題"
    except Exception as e:
        print(f"生成標題時出錯：{e}")
        return "未能生成標題"

def generate_pdf_summary_and_details(pdf_text):
    try:
        summary_prompt = f"請為以下論文文本生成條列式摘要，每點以「-」開頭。請確保每個要點清晰、具體並且包含文檔中的關鍵信息，要求每項要點必須要30 字以上，並用英文說明。\n{pdf_text[:3000]}"
        details_prompt = f"""請為以下論文文本生成更詳細的英文重點整理，每個部分的重點整理應該包括以下內容:
        1. **背景**：簡要介紹文檔的背景，解釋問題的起源和研究的目的。
        2. **方法**：詳細描述文檔中使用的研究方法或分析技巧。
        3. **結果**：呈現研究的主要結果，並詳細描述每個發現的具體細節。
        4. **結論與討論**：分析結果的意涵，並討論研究的貢獻、局限性，以及未來研究的方向。
        每個部分的描述應該詳盡且具體，避免簡單概括，要求更多的上下文和分析，幫助讀者深入理解文檔的精髓。：\n{pdf_text[:3000]}"""

        summary_response = model.generate_content(summary_prompt)
        details_response = model.generate_content(details_prompt)

        # 使用提取函數獲取摘要和詳細內容
        summary = extract_text(summary_response) or "未能生成摘要"
        details = extract_text(details_response) or "未能生成詳細重點整理"

        return summary, details
    except Exception as e:
        print(f"生成摘要與詳細時發生錯誤：{e}")
        return "未能生成摘要。", "未能生成詳細重點整理。"

def ask_llm(prompt):
    try:
        response = call_generative_ai(prompt)
        print(response)
        return response.text.strip() if response else "抱歉，無法生成回答。"
    except Exception as e:
        print(f"LLM 問答失敗: {e}")
        return "抱歉，無法處理您的請求。"

def generate_image_description(image_path):
    # try:
    #     myfile = genai.upload_file(image_path)
    #     prompt = "圖片可能包含各種類型（如圖表、流程圖、模型架構或數據表格），請為圖片生成總共40英文字以內簡短完整的描述。"
    #     result = model.generate_content([myfile, "\n\n", prompt])
    #     return result.text.strip() if result else "未能生成圖片描述"
    # except Exception as e:
    #     print(f"生成圖片描述時出錯：{e}")
    #     return "未能生成圖片描述"
    try:
        # 使用 Base64 將圖片數據傳入 Prompt
        with open(image_path, "rb") as img_file:
            encoded_image = base64.b64encode(img_file.read()).decode("utf-8")
        
        prompt = f"以下是圖片的 Base64 編碼數據，請生成一個簡短的描述（不超過40字）：\n{encoded_image[:500]}"
        response = model.generate_content(prompt)
        
        # 使用提取函數獲取描述文本
        return extract_text(response) or "未能生成圖片描述"
    except Exception as e:
        print(f"生成圖片描述時出錯：{e}")
        return "未能生成圖片描述"

def add_image_and_description_to_slide(prs, image_path, description):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.9), Inches(0.2), Inches(8), Inches(4))

    text_box = slide.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(8), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = description
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.alignment = PP_ALIGN.LEFT

def add_text_slide(prs, title, content, bullet=False, font_size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(font_size + 6)
        paragraph.alignment = PP_ALIGN.LEFT

    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    if bullet:
        for line in content.split("\n"):
            if line.strip():
                paragraph = text_frame.add_paragraph()
                paragraph.text = f"- {line.strip()}"
                paragraph.font.size = Pt(font_size)
                paragraph.space_after = Pt(6)  # 段後間距
    else:
        paragraph = text_frame.add_paragraph()
        paragraph.text = content
        paragraph.font.size = Pt(font_size)
        paragraph.line_spacing = 1.2  # 設置行距
        paragraph.space_after = Pt(12)  # 段後距離

    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT

def split_text_by_points(text, points_per_slide=2):
    points = [line.strip() for line in text.split("\n") if line.strip()]
    return [points[i:i + points_per_slide] for i in range(0, len(points), points_per_slide)]

def adjust_title_font(slide, title):
    font_size = Pt(42) if len(title) <= 50 else Pt(36)
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = font_size

def generate_professor_questions(pdf_text, num_questions=3):
    try:
        prompt = f"基於以下論文內容，生成 {num_questions} 個教授可能問的問題，請用英文：\n{pdf_text[:3000]}"
        response = ask_llm(prompt)
        return response
    except Exception as e:
        print(f"生成教授問題時發生錯誤: {e}")
        return "無法生成教授問題。"
